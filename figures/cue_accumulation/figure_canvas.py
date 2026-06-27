#!/usr/bin/env python3
"""Shared dual-backend canvas for slide-ready, fully-editable task figures.

A figure is described ONCE in a points-based, top-left-origin (y-down) coordinate
system using a small primitive API (line / arrow / rect / polyline / text / dots),
then rendered through two backends so the two output formats are identical:

  * SvgCanvas  -> a clean .svg with real <line>/<rect>/<polyline>/<text> nodes
                  (drop into PowerPoint, "Convert to Shape" -> fully editable),
                  matching the markup conventions of make_deep_leaky_rnn_unrolled.py.
  * PptxCanvas -> a .pptx whose every element is a NATIVE PowerPoint shape
                  (text box, connector, rectangle, freeform polyline) -- recolourable
                  and re-typable with no conversion step.

Design points (pt) map 1:1 to SVG user units; for pptx they convert pt -> EMU.
A 960 x 540 pt canvas equals one 16:9 slide (13.333in x 7.5in).
"""

from __future__ import annotations

# ── shared palette / style (matches existing figures/*.py) ──────────────────
TXT    = "#3f3f3f"   # default text colour
AXIS   = "#7c7c7c"   # axes / arrows / thin rules
LEFT   = "#4C78A8"   # left cue / "rising" accent (cool)
RIGHT  = "#E4794A"   # right cue / "falling" accent (warm)
DELAY  = "#ECECEC"   # silent-delay shading
DELAYE = "#cccccc"   # delay band edge
L1FILL = "#DDEEFF"   # layer-1 box (soft blue, repo pastel)
L1EDGE = "#335577"
L2FILL = "#EEE0FF"   # layer-2 box (soft purple, repo pastel)
L2EDGE = "#553388"
BOXF   = "#FFFFFF"   # generic box fill
RECALL = "#336633"   # recall / decision accent (green, repo convention)
FONT   = "Helvetica Neue, Helvetica, Arial, sans-serif"
FONT_SVG = "'Helvetica Neue', Helvetica, Arial, sans-serif"


def _hex(c: str):
    c = c.lstrip("#")
    return int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)


# ════════════════════════════════════════════════════════════════════════════
# SVG backend
# ════════════════════════════════════════════════════════════════════════════
class SvgCanvas:
    """Accumulates clean SVG primitive nodes and writes a standalone .svg."""

    def __init__(self, width: float, height: float):
        self.W, self.H = width, height
        self._markers: set[str] = set()
        self.body: list[str] = []

    # -- primitives ----------------------------------------------------------
    def line(self, x0, y0, x1, y1, color=AXIS, width=1.6, dash=None):
        d = f' stroke-dasharray="{dash}"' if dash else ""
        self.body.append(
            f'  <line x1="{x0:.1f}" y1="{y0:.1f}" x2="{x1:.1f}" y2="{y1:.1f}" '
            f'stroke="{color}" stroke-width="{width}"{d}/>'
        )

    def arrow(self, x0, y0, x1, y1, color=AXIS, width=1.6):
        mid = self._marker(color)
        self.body.append(
            f'  <line x1="{x0:.1f}" y1="{y0:.1f}" x2="{x1:.1f}" y2="{y1:.1f}" '
            f'stroke="{color}" stroke-width="{width}" marker-end="url(#{mid})"/>'
        )

    def rect(self, x, y, w, h, fill=BOXF, edge=AXIS, width=1.4, radius=0, dash=None):
        f = fill if fill else "none"
        e = edge if edge else "none"
        r = f' rx="{radius}" ry="{radius}"' if radius else ""
        d = f' stroke-dasharray="{dash}"' if dash else ""
        self.body.append(
            f'  <rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
            f'fill="{f}" stroke="{e}" stroke-width="{width}"{r}{d}/>'
        )

    def polyline(self, pts, color=LEFT, width=2.0):
        p = " ".join(f"{x:.1f},{y:.1f}" for x, y in pts)
        self.body.append(
            f'  <polyline points="{p}" fill="none" stroke="{color}" '
            f'stroke-width="{width}" stroke-linejoin="round" stroke-linecap="round"/>'
        )

    def text(self, x, y, s, size=13, color=TXT, anchor="middle",
             weight="normal", italic=False, valign="central"):
        st = ' font-style="italic"' if italic else ""
        self.body.append(
            f'  <text x="{x:.1f}" y="{y:.1f}" font-size="{size}" fill="{color}" '
            f'text-anchor="{anchor}" dominant-baseline="{valign}" '
            f'font-weight="{weight}"{st}>{_esc(s)}</text>'
        )

    def dots(self, x, y, color=AXIS, gap=7, r=2.2):
        for dx in (-gap, 0, gap):
            self.body.append(
                f'  <circle cx="{x+dx:.1f}" cy="{y:.1f}" r="{r}" fill="{color}"/>'
            )

    # -- internals -----------------------------------------------------------
    def _marker(self, color):
        mid = "ah_" + color.lstrip("#")
        if mid not in self._markers:
            self._markers.add(mid)
        return mid

    def _defs(self):
        out = ["  <defs>"]
        for mid in sorted(self._markers):
            color = "#" + mid[len("ah_"):]
            out.append(
                f'    <marker id="{mid}" markerWidth="9" markerHeight="9" refX="7.5" '
                f'refY="3" orient="auto" markerUnits="userSpaceOnUse">'
                f'<path d="M0,0 L8,3 L0,6 Z" fill="{color}"/></marker>'
            )
        out.append("  </defs>")
        return out

    def save(self, path):
        head = (
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{self.W}" '
            f'height="{self.H}" viewBox="0 0 {self.W} {self.H}" '
            f'font-family="{FONT_SVG}">'
        )
        doc = [head] + self._defs() + self.body + ["</svg>"]
        with open(path, "w") as f:
            f.write("\n".join(doc) + "\n")
        print("wrote", path)


def _esc(s: str) -> str:
    return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


# ════════════════════════════════════════════════════════════════════════════
# PPTX backend  (native, editable shapes)
# ════════════════════════════════════════════════════════════════════════════
class PptxCanvas:
    """Mirrors the SvgCanvas API but builds native python-pptx shapes."""

    def __init__(self, width: float, height: float):
        from pptx import Presentation
        from pptx.util import Pt

        self.W, self.H = width, height
        self.prs = Presentation()
        self.prs.slide_width = Pt(width)
        self.prs.slide_height = Pt(height)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

    # -- primitives ----------------------------------------------------------
    def line(self, x0, y0, x1, y1, color=AXIS, width=1.6, dash=None):
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        conn = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Pt(x0), Pt(y0), Pt(x1), Pt(y1)
        )
        conn.line.color.rgb = RGBColor(*_hex(color))
        conn.line.width = Pt(width)
        if dash:
            _set_dash(conn.line)

    def arrow(self, x0, y0, x1, y1, color=AXIS, width=1.6):
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        conn = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Pt(x0), Pt(y0), Pt(x1), Pt(y1)
        )
        conn.line.color.rgb = RGBColor(*_hex(color))
        conn.line.width = Pt(width)
        _set_arrowhead(conn.line)

    def rect(self, x, y, w, h, fill=BOXF, edge=AXIS, width=1.4, radius=0, dash=None):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sp = self.slide.shapes.add_shape(shape_type, Pt(x), Pt(y), Pt(w), Pt(h))
        sp.shadow.inherit = False
        if fill:
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(*_hex(fill))
        else:
            sp.fill.background()
        if edge:
            sp.line.color.rgb = RGBColor(*_hex(edge))
            sp.line.width = Pt(width)
            if dash:
                _set_dash(sp.line)
        else:
            sp.line.fill.background()

    def polyline(self, pts, color=LEFT, width=2.0):
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        builder = self.slide.shapes.build_freeform(Pt(pts[0][0]), Pt(pts[0][1]),
                                                    scale=Pt(1))
        builder.add_line_segments([(Pt(x), Pt(y)) for x, y in pts[1:]], close=False)
        sp = builder.convert_to_shape()
        sp.shadow.inherit = False
        sp.fill.background()
        sp.line.color.rgb = RGBColor(*_hex(color))
        sp.line.width = Pt(width)

    def text(self, x, y, s, size=13, color=TXT, anchor="middle",
             weight="normal", italic=False, valign="central"):
        # (x,y) is the anchor point matching SVG semantics; size in pt.
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml.ns import qn
        # Generous box centred on the anchor; autosize keeps text tight.
        bw, bh = max(40, len(s) * size * 0.62), size * 1.8
        if anchor == "middle":
            left = x - bw / 2
        elif anchor == "start":
            left = x
        else:  # end
            left = x - bw
        top = y - bh / 2
        tb = self.slide.shapes.add_textbox(Pt(left), Pt(top), Pt(bw), Pt(bh))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        # allow the box to grow to fit the text (spaceless autosize)
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        bodyPr.append(bodyPr.makeelement(qn("a:spAutoFit"), {}))
        p = tf.paragraphs[0]
        p.alignment = {"middle": PP_ALIGN.CENTER, "start": PP_ALIGN.LEFT,
                       "end": PP_ALIGN.RIGHT}[anchor]
        run = p.add_run()
        run.text = s
        f = run.font
        f.size = Pt(size)
        f.name = "Helvetica Neue"
        f.bold = weight == "bold"
        f.italic = italic
        f.color.rgb = RGBColor(*_hex(color))

    def dots(self, x, y, color=AXIS, gap=7, r=2.2):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        for dx in (-gap, 0, gap):
            sp = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Pt(x + dx - r), Pt(y - r), Pt(2 * r), Pt(2 * r))
            sp.shadow.inherit = False
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(*_hex(color))
            sp.line.fill.background()

    def save(self, path):
        self.prs.save(path)
        print("wrote", path)


def _set_dash(line):
    from pptx.oxml.ns import qn
    ln = line._get_or_add_ln()
    d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(d)


def _set_arrowhead(line):
    from pptx.oxml.ns import qn
    ln = line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)


# ── convenience: render a figure-builder fn to both backends ────────────────
def render_both(build, width, height, stem, outdir=None):
    """build(canvas) draws the figure; emit <stem>.svg and <stem>.pptx.

    outdir defaults to the directory holding this module (figures/), so the
    scripts write to the right place regardless of the current directory.
    """
    import os
    if outdir is None:
        outdir = os.path.dirname(os.path.abspath(__file__))
    os.makedirs(outdir, exist_ok=True)
    svg = SvgCanvas(width, height)
    build(svg)
    svg.save(os.path.join(outdir, stem + ".svg"))
    pptx = PptxCanvas(width, height)
    build(pptx)
    pptx.save(os.path.join(outdir, stem + ".pptx"))
