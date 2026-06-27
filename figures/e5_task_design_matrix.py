"""
Experiment-5 task-choice design matrix → editable PowerPoint table.

Generates a single-slide ``.pptx`` whose body is a *native* PowerPoint table
(not an image): every cell is editable text, so the table can be copied into a
deck and the labels / ✓ ✗ ? symbols tweaked in PowerPoint itself.

The matrix justifies why **hierarchical cue accumulation** was chosen as the
E5 task: it is the only row that earns ✓ on every trait — it is the only listed
task that requires credit assignment across *both* time and depth at once, is
ablation-isolable, runs in a cheap non-spiking leaky RNN, and is reservoir-proof
(a frozen random lower layer cannot fake the mean-zero motif feature, so depth
credit is genuinely required).

Run:  python figures/e5_task_design_matrix.py
Out:  figures/e5_task_design_matrix.pptx
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Content ──────────────────────────────────────────────────────────────────

TITLE = "Experiment 5 — task choice: why hierarchical cue accumulation"

# Short header text on the slide (long meaning lives in the speaker notes).
COLUMNS = [
    "Temporal\ncredit",
    "Depth\ncredit",
    "Both\nat once",
    "Works\nw/o spikes",
    "Cheap\n(many seeds)",
    "Reservoir-\nproof",
]

# Each row: (task/setup label, [6 cell values], is_chosen)
# Cell values: "Y" = ✓ has it, "N" = ✗ lacks it, "Q" = ? partial/depends.
ROWS = [
    ("Store & recall (vanilla)",        ["N", "N", "N", "Y", "Y", "N"], False),
    ("Store & recall (leaky)",          ["Y", "N", "N", "Y", "Y", "N"], False),
    ("Cue accumulation (vanilla)",      ["N", "N", "N", "Y", "Y", "N"], False),
    ("Cue accumulation (leaky)",        ["Y", "N", "N", "Y", "Y", "N"], False),
    ("Hierarchical cue accum. (leaky)", ["Y", "Y", "Y", "Y", "Y", "Y"], True),
    ("sMNIST (leaky)",                  ["Y", "Q", "N", "Y", "N", "Q"], False),
    ("SHD (spiking LIF/ALIF)",          ["Y", "Q", "N", "N", "Q", "Q"], False),
]

GLYPH = {"Y": "✓", "N": "✗", "Q": "?"}      # ✓ ✗ ?

LEGEND = "✓  has it      ✗  lacks it      ?  partial / depends"

# Per-column meaning → speaker notes (so the slide stays uncluttered).
NOTES = (
    "Why hierarchical cue accumulation (the only all-✓ row):\n"
    "\n"
    "Column meanings:\n"
    "• Temporal credit — the setup requires AND supports non-trivial credit "
    "across time. A vanilla (α=1) tanh RNN has a ~1-step Jacobian, so its "
    "eligibility trace carries nothing and e-prop collapses to d=0 (✗); a leaky "
    "RNN's (1-α) leak carries it forward (✓).\n"
    "• Depth credit — the task needs a LEARNED lower layer (a hierarchy), not "
    "just a trained readout.\n"
    "• Both at once — time and depth credit are entangled and both required "
    "(non-decomposable); this is the actual E5 research question.\n"
    "• Works w/o spikes — runs in a non-spiking leaky/tanh RNN, no "
    "surrogate-gradient SNN needed. SHD is a spiking benchmark, so a non-spiking "
    "net defeats the point (✗).\n"
    "• Cheap (many seeds) — sequences short enough for many-seed significance "
    "testing (16 cosine seeds, n*=8). sMNIST's 784 steps make e-prop too slow (✗).\n"
    "• Reservoir-proof — a frozen RANDOM lower layer fails, so depth credit is "
    "genuinely required and ablate_spatial provably breaks the task. The hierarchical "
    "motifs are mean-zero and equal-energy, so no static/random projection can "
    "separate them — the lower layer MUST learn. Store-and-recall (value is in the "
    "input) and cue accumulation (linearly separable pulses) are reservoir-cheatable "
    "(✗).\n"
    "\n"
    "Takeaway: only hierarchical cue accumulation ticks every box — it is the "
    "decisive simultaneous time-and-depth credit-assignment test."
)

# ── Palette ──────────────────────────────────────────────────────────────────

FILL = {
    "Y": RGBColor(0xC6, 0xEF, 0xCE),   # light green
    "N": RGBColor(0xF2, 0xC8, 0xC8),   # light red
    "Q": RGBColor(0xFF, 0xEB, 0x9C),   # light amber
}
GLYPH_COLOR = {
    "Y": RGBColor(0x1E, 0x6B, 0x2E),   # dark green
    "N": RGBColor(0x9C, 0x1B, 0x1B),   # dark red
    "Q": RGBColor(0x8A, 0x6D, 0x00),   # dark amber
}
HEADER_FILL = RGBColor(0x2F, 0x39, 0x52)   # dark slate
HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
LABEL_FILL = RGBColor(0xF2, 0xF2, 0xF2)    # light grey for task labels
CHOSEN_LABEL_FILL = RGBColor(0xFF, 0xD9, 0x66)   # gold highlight for chosen row
INK = RGBColor(0x22, 0x22, 0x22)


def _set_cell(cell, text, *, bold=False, color=INK, size=12,
              align=PP_ALIGN.CENTER, fill=None):
    """Write text into a table cell with explicit run formatting."""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    font = run.font
    font.name = "Calibri"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color


def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)     # 16:9
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = TITLE
    r.font.name = "Calibri"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = INK

    # Table geometry
    n_rows = len(ROWS) + 1
    n_cols = len(COLUMNS) + 1
    left, top = Inches(0.5), Inches(1.25)
    width, height = Inches(12.33), Inches(5.3)
    graphic = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic.table

    # Column widths: wide label column, even trait columns.
    label_w = Inches(3.2)
    trait_w = Emu(int((width - label_w) / len(COLUMNS)))
    table.columns[0].width = label_w
    for c in range(1, n_cols):
        table.columns[c].width = trait_w

    # Header row
    _set_cell(table.cell(0, 0), "Task (setup)", bold=True, color=HEADER_TEXT,
              size=13, align=PP_ALIGN.LEFT, fill=HEADER_FILL)
    for j, name in enumerate(COLUMNS, start=1):
        _set_cell(table.cell(0, j), name, bold=True, color=HEADER_TEXT,
                  size=12, fill=HEADER_FILL)

    # Data rows
    for i, (label, values, chosen) in enumerate(ROWS, start=1):
        _set_cell(table.cell(i, 0), label, bold=chosen, color=INK, size=12,
                  align=PP_ALIGN.LEFT,
                  fill=CHOSEN_LABEL_FILL if chosen else LABEL_FILL)
        for j, v in enumerate(values, start=1):
            _set_cell(table.cell(i, j), GLYPH[v], bold=True,
                      color=GLYPH_COLOR[v], size=16, fill=FILL[v])

    # Legend
    lb = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5))
    lp = lb.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = LEGEND
    lr.font.name = "Calibri"
    lr.font.size = Pt(13)
    lr.font.color.rgb = INK

    # Speaker notes
    slide.notes_slide.notes_text_frame.text = NOTES

    prs.save(path)
    return path


if __name__ == "__main__":
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "e5_task_design_matrix.pptx")
    build(out)
    print(f"wrote {out}")
